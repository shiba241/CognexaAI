import subprocess
import sys
subprocess.check_call([sys.executable, "-m", "pip", "install", "python-pptx", "--quiet"])

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
import os

# Paths
artifact_dir = r"C:\Users\shiba\.gemini\antigravity\brain\bbb2d42e-362c-4b00-a7b5-51bb6704dcf4"
output_path = r"c:\Users\shiba\OneDrive\Desktop\Cognexa AI-Autonomous-Multi-Agent-Research-System\ACE_Professional_Presentation.pptx"

# Find image files
def find_image(prefix):
    for f in os.listdir(artifact_dir):
        if f.startswith(prefix) and f.endswith('.png'):
            return os.path.join(artifact_dir, f)
    return None

# Create presentation (Widescreen 16:9)
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Color constants
DARK_BG = RGBColor(0x0B, 0x0E, 0x1A)
CYAN = RGBColor(0x00, 0xD4, 0xFF)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xB0, 0xB8, 0xC8)
TEAL = RGBColor(0x00, 0xBF, 0xA6)

def set_slide_bg(slide, color=DARK_BG):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_title(slide, text, left=0.5, top=0.3, width=12.3, height=0.9, size=36, color=CYAN, bold=True):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = "Segoe UI"
    p.alignment = PP_ALIGN.LEFT
    return txBox

def add_subtitle(slide, text, left=0.5, top=1.2, width=12.3, height=0.6, size=18, color=LIGHT_GRAY):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.name = "Segoe UI"
    p.alignment = PP_ALIGN.LEFT
    return txBox

def add_bullet_text(slide, bullets, left=0.5, top=2.0, width=5.5, height=4.5, size=16, color=WHITE):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, (title, desc) in enumerate(bullets):
        if i > 0:
            p = tf.add_paragraph()
            p.text = ""
            p.space_before = Pt(6)
        # Title bullet
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.text = f"▸ {title}"
        p.font.size = Pt(size + 2)
        p.font.color.rgb = CYAN
        p.font.bold = True
        p.font.name = "Segoe UI"
        p.space_before = Pt(12) if i > 0 else Pt(0)
        # Description
        p2 = tf.add_paragraph()
        p2.text = f"   {desc}"
        p2.font.size = Pt(size - 1)
        p2.font.color.rgb = LIGHT_GRAY
        p2.font.name = "Segoe UI"
        p2.space_before = Pt(4)
    return txBox

def add_image(slide, img_prefix, left=6.5, top=1.8, width=6.3, height=5.0):
    img_path = find_image(img_prefix)
    if img_path:
        slide.shapes.add_picture(img_path, Inches(left), Inches(top), Inches(width), Inches(height))

def add_full_image(slide, img_prefix, left=0.5, top=1.8, width=12.3, height=5.2):
    img_path = find_image(img_prefix)
    if img_path:
        slide.shapes.add_picture(img_path, Inches(left), Inches(top), Inches(width), Inches(height))

def add_speaker_note(slide, text):
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = text

def add_slide_number(slide, num, total=13):
    txBox = slide.shapes.add_textbox(Inches(12.0), Inches(7.0), Inches(1.0), Inches(0.4))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = f"{num} / {total}"
    p.font.size = Pt(10)
    p.font.color.rgb = RGBColor(0x60, 0x68, 0x78)
    p.font.name = "Segoe UI"
    p.alignment = PP_ALIGN.RIGHT

# ==================== SLIDE 1: TITLE ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
set_slide_bg(slide)
img_path = find_image("slide1_title_background")
if img_path:
    slide.shapes.add_picture(img_path, Inches(0), Inches(0), Inches(13.333), Inches(7.5))

# Title text overlay
txBox = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(11), Inches(1.5))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Autonomous Cognitive Engine (ACE)"
p.font.size = Pt(44)
p.font.color.rgb = WHITE
p.font.bold = True
p.font.name = "Segoe UI"
p.alignment = PP_ALIGN.CENTER

p2 = tf.add_paragraph()
p2.text = "For Deep Research"
p2.font.size = Pt(36)
p2.font.color.rgb = CYAN
p2.font.bold = True
p2.font.name = "Segoe UI"
p2.alignment = PP_ALIGN.CENTER

# Subtitle
txBox2 = slide.shapes.add_textbox(Inches(1), Inches(3.5), Inches(11), Inches(0.8))
tf2 = txBox2.text_frame
tf2.word_wrap = True
p3 = tf2.paragraphs[0]
p3.text = "AI-Powered Multi-Agent System That Plans, Researches, Analyzes & Delivers Reports Autonomously"
p3.font.size = Pt(18)
p3.font.color.rgb = LIGHT_GRAY
p3.font.name = "Segoe UI"
p3.alignment = PP_ALIGN.CENTER

# Author info
txBox3 = slide.shapes.add_textbox(Inches(1), Inches(5.0), Inches(11), Inches(1.5))
tf3 = txBox3.text_frame
tf3.word_wrap = True
p4 = tf3.paragraphs[0]
p4.text = "Presented by: Shiba Prasad Sahu"
p4.font.size = Pt(20)
p4.font.color.rgb = WHITE
p4.font.name = "Segoe UI"
p4.alignment = PP_ALIGN.CENTER

p5 = tf3.add_paragraph()
p5.text = "Program: PMIS Internship Program — Infosys Limited"
p5.font.size = Pt(16)
p5.font.color.rgb = LIGHT_GRAY
p5.font.name = "Segoe UI"
p5.alignment = PP_ALIGN.CENTER

add_speaker_note(slide, "Welcome everyone. Today I'll present the Autonomous Cognitive Engine — an AI system that goes beyond simple chatbots to autonomously plan, research, and deliver professional reports.")

# ==================== SLIDE 2: CORE CHALLENGE ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_title(slide, "The Core Challenge: Beyond Standard LLMs")
add_subtitle(slide, "Why current AI models fail at deep, multi-step research tasks")
add_bullet_text(slide, [
    ("They Forget Mid-Task", "LLMs have a fixed context window — during long research, they lose early instructions and start hallucinating facts"),
    ("No Planning or Delegation", "A single LLM can't break a complex problem into subtasks or assign specialized roles — it works like one person doing everything alone"),
    ("No Self-Correction", "Standard LLMs deliver output in one shot — no review, no quality check, no retry if the result is poor"),
])
add_image(slide, "slide2_core_challenges")
add_slide_number(slide, 2)
add_speaker_note(slide, "To understand the value of ACE, we must look at the limitations of current AI. Modern LLMs are powerful but they're short-term thinkers — they forget, work alone, and never check their own work.")

# ==================== SLIDE 3: ACE SOLUTION ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_title(slide, "The ACE Solution: A Deep Cognitive Framework")
add_subtitle(slide, "From chatbot to intelligent coordinator")
add_bullet_text(slide, [
    ("Works Like a Manager, Not a Chatbot", "ACE acts as a central coordinator — it reads your goal, creates a strict step-by-step plan, and controls the entire execution flow"),
    ("Never Forgets — Uses a Virtual File System", "Every agent saves its work into a shared memory (VFS), so no information is lost between steps"),
    ("Delegates Tasks to Specialist Agents", "Instead of one AI doing everything, ACE assigns each task to a focused agent — search, analysis, examples — each expert in its own job"),
])
add_image(slide, "slide3_ace_solution")
add_slide_number(slide, 3)
add_speaker_note(slide, "ACE shifts AI from a chatbot to a project manager — it plans, remembers using a virtual file system, and delegates to the right specialists automatically.")

# ==================== SLIDE 4: TECH STACK ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_title(slide, "Enterprise-Grade Technology Stack")
add_subtitle(slide, "Built for scale, speed, and observability")
add_bullet_text(slide, [
    ("Python 3.11+ with LangGraph & LangChain", "Python handles core logic, LangGraph manages multi-step workflows with shared state, LangChain connects to LLM APIs"),
    ("Advanced LLMs for Decision-Making", "LLMs (Gemini/Claude) act as the thinking engine — they plan, reason, and act using the ReAct pattern"),
    ("LangSmith for Real-Time Monitoring", "Tracks every step, every agent call, and every decision in real time — making debugging and evaluation easy"),
])
add_image(slide, "slide4_tech_stack")
add_slide_number(slide, 4)
add_speaker_note(slide, "Our stack is built for scale — Python and LangGraph handle the workflow, advanced LLMs drive the thinking, and LangSmith lets us see exactly what's happening inside at every step.")

# ==================== SLIDE 5: WHY LANGGRAPH ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_title(slide, "Architectural Shift: Why LangGraph?")
add_subtitle(slide, "Moving from linear pipelines to cyclic, stateful graphs")
add_bullet_text(slide, [
    ("Supports Loops & Retries", "Unlike basic LangChain pipelines that run in a straight line, LangGraph allows the AI to go back, retry failed steps, and repeat until quality is met"),
    ("Every Step is a Separate Node", "Each action — planning, searching, analyzing, writing — is a modular building block connected by edges, easy to control and debug"),
    ("Keeps Memory Alive Across All Steps", "LangGraph maintains shared state throughout the entire process — no step loses track of what previous steps did"),
])
add_image(slide, "slide5_langgraph_comparison")
add_slide_number(slide, 5)
add_speaker_note(slide, "Research isn't a straight line — sometimes you need to loop back and retry. LangGraph lets our agents do exactly that, while keeping full memory at every step.")

# ==================== SLIDE 6: ARCHITECTURE DIAGRAM ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_title(slide, "System Architecture: LangGraph StateGraph")
add_subtitle(slide, "End-to-end autonomous research pipeline")
# Full width image for architecture
add_full_image(slide, "slide6_system_architecture", left=0.3, top=1.8, width=8.0, height=5.2)
# Description on right
txBox = slide.shapes.add_textbox(Inches(8.5), Inches(2.0), Inches(4.3), Inches(5.0))
tf = txBox.text_frame
tf.word_wrap = True
items = [
    ("1.", "User enters research topic"),
    ("2.", "Planner breaks it into TODO tasks"),
    ("3.", "3 Workers execute in sequence"),
    ("4.", "Results saved to Virtual File System"),
    ("5.", "Synthesizer merges into full report"),
    ("6.", "Quality Gate scores (0-10)"),
    ("7.", "Score < 7 → auto-retry"),
    ("8.", "Score ≥ 7 → deliver final report"),
]
for i, (num, text) in enumerate(items):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.text = f"  {num}  {text}"
    p.font.size = Pt(14)
    p.font.color.rgb = LIGHT_GRAY if i % 2 == 0 else WHITE
    p.font.name = "Segoe UI"
    p.space_before = Pt(8)
add_slide_number(slide, 6)
add_speaker_note(slide, "The entire system runs on a LangGraph StateGraph — a shared brain that keeps all agents in sync. Every agent reads from it and writes back to it, so the system stays perfectly synchronized.")

# ==================== SLIDE 7: REACT LOOP ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_title(slide, "The Core Engine: ReAct Execution Loop")
add_subtitle(slide, "Every agent follows a disciplined Reason → Act → Observe cycle")
add_bullet_text(slide, [
    ("Think Before Acting", "The agent first reasons about the current situation and decides the best next step"),
    ("Act by Calling the Right Tool", "Based on its reasoning, it picks the most suitable action — searching, analyzing, or writing"),
    ("Observe the Result Before Moving On", "After every action, it checks the outcome — if something fails, it adapts instead of making up an answer"),
])
add_image(slide, "slide7_react_loop")
add_slide_number(slide, 7)
add_speaker_note(slide, "At the heart of every agent is the ReAct loop. It thinks first, takes action, then checks the result — if something fails, it adapts instead of guessing.")

# ==================== SLIDE 8: TASK PLANNING ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_title(slide, "Intelligent Task Planning (write_todos)")
add_subtitle(slide, "Forced planning before execution — no shortcuts allowed")
add_bullet_text(slide, [
    ("Plans First, Executes Second", "ACE analyzes the research goal and breaks it down into smaller, clear sub-tasks — like a manager creating a to-do list before assigning work"),
    ("write_todos Tool Forces Structure", "The AI is required to use the write_todos tool — it can't skip planning, ensuring every task gets a proper action plan"),
    ("Every Task is Tracked to Completion", "All sub-tasks are saved into the LangGraph state and tracked one by one — nothing gets missed"),
])
add_image(slide, "slide8_task_planning")
add_slide_number(slide, 8)
add_speaker_note(slide, "ACE is forced to plan before it acts — it breaks down big goals into trackable sub-tasks using write_todos, ensuring nothing is skipped and every step is accounted for.")

# ==================== SLIDE 9: VFS ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_title(slide, "Context Offloading: The Virtual File System")
add_subtitle(slide, "How agents save, share, and retrieve information without memory overload")
# Full image for VFS diagram
add_full_image(slide, "slide9_vfs_diagram", left=0.3, top=1.8, width=8.0, height=5.0)
# Description on right
txBox = slide.shapes.add_textbox(Inches(8.5), Inches(2.2), Inches(4.3), Inches(4.5))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Instead of holding everything in memory at once, agents write their work into virtual files and read back only what they need."
p.font.size = Pt(15)
p.font.color.rgb = LIGHT_GRAY
p.font.name = "Segoe UI"
p2 = tf.add_paragraph()
p2.text = ""
p3 = tf.add_paragraph()
p3.text = "▸ Tools: write_file, read_file, ls, edit_file"
p3.font.size = Pt(14)
p3.font.color.rgb = CYAN
p3.font.name = "Segoe UI"
p3.space_before = Pt(12)
p4 = tf.add_paragraph()
p4.text = "▸ Storage: Python dictionary in state"
p4.font.size = Pt(14)
p4.font.color.rgb = CYAN
p4.font.name = "Segoe UI"
p4.space_before = Pt(8)
p5 = tf.add_paragraph()
p5.text = "▸ Speed: In-memory, zero disk I/O"
p5.font.size = Pt(14)
p5.font.color.rgb = CYAN
p5.font.name = "Segoe UI"
p5.space_before = Pt(8)
p6 = tf.add_paragraph()
p6.text = "▸ Benefit: Breaks through LLM context limits and saves token costs"
p6.font.size = Pt(14)
p6.font.color.rgb = CYAN
p6.font.name = "Segoe UI"
p6.space_before = Pt(8)
add_slide_number(slide, 9)
add_speaker_note(slide, "The VFS is the AI's notebook — agents write findings, free up memory, and read back only what they need. This saves cost and lets the system handle far more data than a standard LLM can.")

# ==================== SLIDE 10: MODULAR INTELLIGENCE ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_title(slide, "Modular Intelligence: Sub-Agent Delegation")
add_subtitle(slide, "Supervisor assigns tasks — each agent works in isolation")
add_bullet_text(slide, [
    ("Supervisor Acts as Manager", "The main agent doesn't do the heavy lifting — it assigns specific tasks to specialized sub-agents"),
    ("Isolated Contexts Prevent Interference", "Each agent works in its own space with only the tools it needs — no data contamination between agents"),
    ("Focused Results, No Overload", "Specialization means higher quality outputs — the main engine stays light and the agents stay focused"),
], width=5.0)
add_image(slide, "slide10_modular_agents")
add_slide_number(slide, 10)
add_speaker_note(slide, "The main supervisor doesn't do all the heavy lifting itself — it delegates to specialized sub-agents. Each works in isolation, stays focused, and the main engine never gets overloaded.")

# ==================== SLIDE 11: QUALITY ASSURANCE ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_title(slide, "Quality Assurance & Evaluation Metrics")
add_subtitle(slide, "Rigorous testing verified through LangSmith traces")
add_bullet_text(slide, [
    ("Task Planning — 80%+ Accuracy", "The Planner successfully broke down complex research goals into clear, actionable sub-tasks in over 80% of test cases"),
    ("VFS Usage — 80%+ Autonomous Adoption", "Agents correctly saved and retrieved data using the Virtual File System in more than 80% of multi-step scenarios"),
    ("Output Quality — LLM-as-a-Judge", "Every final report is scored by the AI itself on a 0–10 scale. Score below 7 triggers automatic retry"),
])
add_image(slide, "slide11_quality_metrics")
add_slide_number(slide, 11)
add_speaker_note(slide, "The system was tested rigorously — 80%+ planning accuracy, 80%+ VFS adoption, and automated quality scoring with auto-retry for reports below standard.")

# ==================== SLIDE 12: CHALLENGES ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_title(slide, "Technical Challenges & Key Learnings")
add_subtitle(slide, "Real engineering problems we solved")
add_bullet_text(slide, [
    ("Controlling Unpredictable AI Behavior", "LLMs can loop endlessly or call tools incorrectly — we solved this by designing strict system prompts and enforcing Pydantic schemas to keep agents disciplined"),
    ("Managing Shared Memory Without Data Loss", "Multiple agents writing to the same state can overwrite each other's work — we learned to carefully synchronize state updates so every agent's output stays safe"),
])
add_image(slide, "slide12_challenges")
add_slide_number(slide, 12)
add_speaker_note(slide, "We encountered significant challenges in making unpredictable LLMs follow strict workflows. Strict prompts, Pydantic schemas, and careful state management were key solutions.")

# ==================== SLIDE 13: LIVE DEMO + THANK YOU ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
img_path = find_image("slide13_thank_you")
if img_path:
    slide.shapes.add_picture(img_path, Inches(0), Inches(0), Inches(13.333), Inches(7.5))

txBox = slide.shapes.add_textbox(Inches(1), Inches(1.0), Inches(11), Inches(1.2))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Live Demonstration"
p.font.size = Pt(40)
p.font.color.rgb = CYAN
p.font.bold = True
p.font.name = "Segoe UI"
p.alignment = PP_ALIGN.CENTER

txBox2 = slide.shapes.add_textbox(Inches(2), Inches(2.5), Inches(9), Inches(1.5))
tf2 = txBox2.text_frame
tf2.word_wrap = True
p2 = tf2.paragraphs[0]
p2.text = "🔗  https://cognexaa.streamlit.app/"
p2.font.size = Pt(24)
p2.font.color.rgb = WHITE
p2.font.name = "Segoe UI"
p2.alignment = PP_ALIGN.CENTER

txBox3 = slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(11), Inches(1.5))
tf3 = txBox3.text_frame
tf3.word_wrap = True
p3 = tf3.paragraphs[0]
p3.text = "Thank You"
p3.font.size = Pt(48)
p3.font.color.rgb = WHITE
p3.font.bold = True
p3.font.name = "Segoe UI"
p3.alignment = PP_ALIGN.CENTER

p4 = tf3.add_paragraph()
p4.text = "Shiba Prasad Sahu  |  PMIS Internship Program — Infosys Limited"
p4.font.size = Pt(16)
p4.font.color.rgb = LIGHT_GRAY
p4.font.name = "Segoe UI"
p4.alignment = PP_ALIGN.CENTER
p4.space_before = Pt(12)

add_speaker_note(slide, "That concludes the overview. Let's now switch to a live demo where ACE will research a complex topic in real time while we monitor it through LangSmith. Thank you.")

# Save
prs.save(output_path)
print(f"PPT saved to: {output_path}")
print(f"Total slides: {len(prs.slides)}")
