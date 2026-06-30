import subprocess
import sys
subprocess.check_call([sys.executable, "-m", "pip", "install", "python-pptx", "--quiet"])

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
import os

# Paths
artifact_dir = r"C:\Users\shiba\.gemini\antigravity\brain\bbb2d42e-362c-4b00-a7b5-51bb6704dcf4"
output_path = r"c:\Users\shiba\OneDrive\Desktop\Cognexa AI-Autonomous-Multi-Agent-Research-System\CognexaAI_Team_Presentation.pptx"

# Find image files
def find_image(prefix):
    for f in os.listdir(artifact_dir):
        if f.startswith(prefix) and f.endswith('.png'):
            return os.path.join(artifact_dir, f)
    return None

# Create presentation (Widescreen 16:9)
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Color constants
DARK_BG = RGBColor(0x0B, 0x0E, 0x1A)
CYAN = RGBColor(0x00, 0xD4, 0xFF)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xB0, 0xB8, 0xC8)
TEAL = RGBColor(0x00, 0xBF, 0xA6)

def set_slide_bg(slide, color=DARK_BG):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_title(slide, text, left=0.5, top=0.3, width=12.3, height=0.9, size=36, color=CYAN, bold=True):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = "Segoe UI"
    p.alignment = PP_ALIGN.LEFT
    return txBox

def add_subtitle(slide, text, left=0.5, top=1.2, width=12.3, height=0.6, size=18, color=LIGHT_GRAY):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.name = "Segoe UI"
    p.alignment = PP_ALIGN.LEFT
    return txBox

def add_bullet_text(slide, bullets, left=0.5, top=2.0, width=5.5, height=4.5, size=18, color=WHITE):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, title in enumerate(bullets):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.text = f"▸ {title}"
        p.font.size = Pt(size)
        p.font.color.rgb = WHITE
        p.font.name = "Segoe UI"
        p.space_before = Pt(16)
    return txBox

def add_image(slide, img_prefix, left=6.5, top=1.8, width=6.3, height=5.0):
    img_path = find_image(img_prefix)
    if img_path:
        slide.shapes.add_picture(img_path, Inches(left), Inches(top), Inches(width), Inches(height))

def add_slide_number(slide, num, total=13):
    txBox = slide.shapes.add_textbox(Inches(12.0), Inches(7.0), Inches(1.0), Inches(0.4))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = f"{num} / {total}"
    p.font.size = Pt(10)
    p.font.color.rgb = RGBColor(0x60, 0x68, 0x78)
    p.font.name = "Segoe UI"
    p.alignment = PP_ALIGN.RIGHT

# ==================== SLIDE 1: TITLE ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
img_path = find_image("human_team_bg")
if img_path:
    slide.shapes.add_picture(img_path, Inches(0), Inches(0), Inches(13.333), Inches(7.5))

txBox = slide.shapes.add_textbox(Inches(1), Inches(1.0), Inches(11), Inches(1.5))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "CognexaAI"
p.font.size = Pt(54)
p.font.color.rgb = WHITE
p.font.bold = True
p.font.name = "Segoe UI"

p2 = tf.add_paragraph()
p2.text = "Autonomous Multi-Agent Research System"
p2.font.size = Pt(28)
p2.font.color.rgb = TEAL
p2.font.bold = False
p2.font.name = "Segoe UI"
p2.space_before = Pt(10)

txBox3 = slide.shapes.add_textbox(Inches(1), Inches(3.5), Inches(5), Inches(3.0))
tf3 = txBox3.text_frame
tf3.word_wrap = True
team = [
    "Team: Subhashree Pattnaik (Team Lead)",
    "Shiba Prasad Sahu",
    "Laxmiranjan Sahu",
    "Kanhu Charan Sahoo",
    "Vetcha Shyam",
    "SmrutiRanjan Jena",
    "Chirasmita Prusty"
]
for i, member in enumerate(team):
    p = tf3.add_paragraph() if i > 0 else tf3.paragraphs[0]
    p.text = member
    p.font.size = Pt(18) if i > 0 else Pt(20)
    p.font.color.rgb = WHITE
    p.font.bold = True if i == 0 else False
    p.font.name = "Segoe UI"

txBox4 = slide.shapes.add_textbox(Inches(7), Inches(4.5), Inches(5), Inches(1.0))
tf4 = txBox4.text_frame
p4 = tf4.paragraphs[0]
p4.text = "Infosys × PMIS"
p4.font.size = Pt(28)
p4.font.color.rgb = TEAL
p4.font.bold = True
p4.font.name = "Segoe UI"

# ==================== SLIDE 2: AGENDA ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_title(slide, "Agenda")
add_subtitle(slide, "Presentation Roadmap")
agenda_left = [
    "1. Problem Statement & Solution Overview",
    "2. Technology Stack & Design",
    "3. System Architecture",
    "4. Pipeline Workflow",
    "5. Virtual File System",
    "6. Sub Agents Overview"
]
agenda_right = [
    "7. Evaluation & Quality Gate",
    "8. Security & Privacy",
    "9. Engineering Challenges",
    "10. Future Works Roadmap",
    "11. Live Demo",
    "12. Questions & Answers"
]
txBoxL = slide.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(5.5), Inches(4.5))
tfL = txBoxL.text_frame
for i, item in enumerate(agenda_left):
    p = tfL.add_paragraph() if i > 0 else tfL.paragraphs[0]
    p.text = item
    p.font.size = Pt(20)
    p.font.color.rgb = WHITE
    p.font.name = "Segoe UI"
    p.space_before = Pt(18)

txBoxR = slide.shapes.add_textbox(Inches(7.0), Inches(2.0), Inches(5.5), Inches(4.5))
tfR = txBoxR.text_frame
for i, item in enumerate(agenda_right):
    p = tfR.add_paragraph() if i > 0 else tfR.paragraphs[0]
    p.text = item
    p.font.size = Pt(20)
    p.font.color.rgb = WHITE
    p.font.name = "Segoe UI"
    p.space_before = Pt(18)
add_slide_number(slide, 2)

# ==================== SLIDE 3: PROBLEM & SOLUTION ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_title(slide, "Problem Statement & Solution")
add_subtitle(slide, "Why we built CognexaAI")
add_bullet_text(slide, [
    "Standard AI forgets facts and guesses answers without live data.",
    "Too many tasks overload standard AI, leading to poor output.",
    "CognexaAI fixes this by splitting work into specialized steps.",
    "Specialized agents focus on one task: searching, analyzing, or writing."
], width=6.0)
add_image(slide, "human_problem_solution", left=7.0, top=1.8, width=5.5, height=5.0)
add_slide_number(slide, 3)

# ==================== SLIDE 4: TECH STACK ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_title(slide, "Technology Stack")
add_subtitle(slide, "Tools and frameworks powering the system")
add_bullet_text(slide, [
    "Python 3.11+: Handles all core system logic.",
    "LangGraph: Manages the step-by-step workflow smoothly.",
    "Gemini & Groq: The smart AI brains making decisions.",
    "Tavily API: Gathers real-time information from the web.",
    "Streamlit: Provides a simple, interactive user website."
], width=6.0)
add_image(slide, "human_tech_stack", left=7.0, top=1.8, width=5.5, height=5.0)
add_slide_number(slide, 4)

# ==================== SLIDE 5: ARCHITECTURE ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_title(slide, "System Architecture")
add_subtitle(slide, "How the pieces connect")
add_bullet_text(slide, [
    "Planner Node: Creates a clear list of tasks to do.",
    "Workers Node: Searches the web, analyzes data, and finds examples.",
    "Synthesizer Node: Combines everything into a final report.",
    "Quality Gate Node: Checks the report and retries if it is bad.",
    "The system runs in a loop until the quality is perfect."
], width=6.0)
add_image(slide, "human_workflow", left=7.0, top=1.8, width=5.5, height=5.0)
add_slide_number(slide, 5)

# ==================== SLIDE 6: WORKFLOW ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_title(slide, "Pipeline Workflow")
add_subtitle(slide, "Step-by-step Execution")
add_bullet_text(slide, [
    "Step 1: User enters a research topic in the app.",
    "Step 2: AI automatically creates a 3-step action plan.",
    "Step 3: Agents execute the plan to gather information.",
    "Step 4: AI writes and formats the final document.",
    "Step 5: The final report is shown and can be downloaded."
], width=11.0)
add_slide_number(slide, 6)

# ==================== SLIDE 7: VFS ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_title(slide, "Virtual File System")
add_subtitle(slide, "In-Memory Workspace")
add_bullet_text(slide, [
    "Agents save their work in virtual memory, not on the hard drive.",
    "This avoids messy files and keeps the system lightning fast.",
    "Each agent can easily read what the previous agent wrote.",
    "Keeps the AI focused without overloading its memory."
], width=11.0)
add_slide_number(slide, 7)

# ==================== SLIDE 8: SUB AGENTS ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_title(slide, "Sub Agents")
add_subtitle(slide, "3 Specialist Workers in Detail")
add_bullet_text(slide, [
    "Search Agent: Finds raw web articles and information.",
    "Analysis Agent: Reads articles and summarizes the key points.",
    "Examples Agent: Finds real company case studies and results.",
    "Working separately gives much better results than doing it all at once."
], width=11.0)
add_slide_number(slide, 8)

# ==================== SLIDE 9: EVALUATION ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_title(slide, "Evaluation")
add_subtitle(slide, "Quality Gate & Self-Correction")
add_bullet_text(slide, [
    "The AI scores its own final report on a scale from 0 to 10.",
    "If the score is low, the AI automatically goes back to fix it.",
    "It has a maximum of 2 retries to prevent endless loops.",
    "This ensures the user always receives high-quality output."
], width=11.0)
add_slide_number(slide, 9)

# ==================== SLIDE 10: SECURITY ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_title(slide, "Security")
add_subtitle(slide, "API Key Isolation & Sanitization")
add_bullet_text(slide, [
    "API keys are kept strictly hidden and never uploaded online.",
    "Text is safely cleaned before creating the PDF to prevent errors.",
    "The virtual workspace ensures no real computer files are ever deleted.",
    "The system is completely isolated and safe to use."
], width=11.0)
add_slide_number(slide, 10)

# ==================== SLIDE 11: CHALLENGES ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_title(slide, "Engineering Challenges")
add_subtitle(slide, "Problems we faced and solved")
add_bullet_text(slide, [
    "Challenge: Free AI hit limits quickly. Solution: Added Groq as backup.",
    "Challenge: Data got lost between steps. Solution: Fixed memory saving logic.",
    "Challenge: PDF creation crashed on symbols. Solution: Built a text cleaner.",
    "Challenge: AI skipped steps. Solution: Forced strict planning rules."
], width=11.0)
add_slide_number(slide, 11)

# ==================== SLIDE 12: FUTURE WORKS ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_title(slide, "Future Works")
add_subtitle(slide, "Extension Roadmap")
add_bullet_text(slide, [
    "Run all agents at the same time for faster results.",
    "Allow users to upload their own PDFs for the AI to read.",
    "Automatically add proper academic citations to the report.",
    "Save previous chats so users can resume their research later."
], width=11.0)
add_slide_number(slide, 12)

# ==================== SLIDE 13: THANK YOU ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

txBox = slide.shapes.add_textbox(Inches(1), Inches(2.0), Inches(11), Inches(1.2))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "Thank You!"
p.font.size = Pt(64)
p.font.color.rgb = CYAN
p.font.bold = True
p.font.name = "Segoe UI"
p.alignment = PP_ALIGN.CENTER

txBox2 = slide.shapes.add_textbox(Inches(1), Inches(3.5), Inches(11), Inches(1.0))
tf2 = txBox2.text_frame
p2 = tf2.paragraphs[0]
p2.text = "We are ready for your questions — every team member can answer"
p2.font.size = Pt(24)
p2.font.color.rgb = WHITE
p2.font.name = "Segoe UI"
p2.alignment = PP_ALIGN.CENTER

txBox3 = slide.shapes.add_textbox(Inches(1), Inches(5.0), Inches(11), Inches(1.5))
tf3 = txBox3.text_frame
p3 = tf3.paragraphs[0]
p3.text = "Live App: cognexaa.streamlit.app"
p3.font.size = Pt(28)
p3.font.color.rgb = TEAL
p3.font.bold = True
p3.font.name = "Segoe UI"
p3.alignment = PP_ALIGN.CENTER

prs.save(output_path)
print(f"PPT saved to: {output_path}")
print(f"Total slides: {len(prs.slides)}")
